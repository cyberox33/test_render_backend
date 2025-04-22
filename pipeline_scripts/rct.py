import os
import pandas as pd
from supabase import create_client, Client
from google import generativeai
import plotly.express as px
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from datetime import datetime, timedelta
import io
import textwrap

# ---- SETUP ----

# Supabase
SUPABASE_URL = "https://rahcdnyvtrhdauixoyzp.supabase.co"
SUPABASE_KEY = "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJpc3MiOiJzdXBhYmFzZSIsInJlZiI6InJhaGNkbnl2dHJoZGF1aXhveXpwIiwicm9sZSI6ImFub24iLCJpYXQiOjE3NDA0MDU3NDAsImV4cCI6MjA1NTk4MTc0MH0.knl6CJKkLahwonX3CTHMPMMekOpgn1gXejfOJ-3UUaA"
supabase: Client = create_client(SUPABASE_URL, SUPABASE_KEY)

# Google GenAI
generativeai.configure(api_key="AIzaSyCu7LPLrOExMye01c2PA9BZodMpDn-KE-A")
llm = generativeai.GenerativeModel("gemma-3-27b-it")

session_id = "152f055d-122a-4793-8305-f9cf4997b75e"

# ---- RETRIEVE DATA ----

response = supabase.table("followup_responses").select("*").eq("session_id", session_id).execute()
followups = response.data

question_ids = [f['question_id'] for f in followups]
questions_data = supabase.table("questions").select("id, question, additional_fields").in_("id", question_ids).execute().data
question_map = {q['id']: q for q in questions_data}

# ---- STRUCTURE COMBINED DATA ----

fragment_groups = {}
for f in followups:
    q_data = question_map.get(f['question_id'], {})
    fragment = f['Fragment']
    combined = {
        "question_id": f['question_id'],
        "question": q_data.get("question", f["question"]),
        "additional_fields": q_data.get("additional_fields", {}),
        "fragment": fragment,
        "answer": f['answer']
    }
    fragment_groups.setdefault(fragment, []).append(combined)

# ---- PROMPT FUNCTIONS ----

def prompt_generate_recommendation(q):
    base_prompt = f"""Given the question and user answer below, determine whether a recommendation is needed and if so, generate it using the structured schema.

Question: {q['question']}
Answer: {q['answer']}
Additional Info: {q['additional_fields']}

Follow the format of:
- Step #: Description
- man_hour_total
- estimate_man_hours
- resources_skills_required

Only generate if needed. Respond with structured JSON or explain why none is needed."""
    result = llm.generate_content(base_prompt)
    return result.text

def prompt_convert_fragment_A(q):
    base_prompt = f"""This question belongs to Fragment A.

Try to convert its recommendation and additional info to match the schema of B/C (steps, man_hour_total, skills, etc.).

Question: {q['question']}
Answer: {q['answer']}
Additional Fields: {q['additional_fields']}

Respond in structured JSON with steps, man_hour_total, estimate_man_hours, resources_skills_required."""
    result = llm.generate_content(base_prompt)
    return result.text

def prompt_summarize_fragment(fragment, questions):
    base_prompt = f"""You're analyzing Fragment {fragment} across multiple questions.

Summarize the roadmap step-by-step by merging all their structured recommendations into a coherent progression.

Questions and their structured recommendations:
{questions}

Return a structured summary with step number, description, and grouped man-hours per phase if available."""
    result = llm.generate_content(base_prompt)
    return result.text

# ---- NORMALIZATION + PROCESSING ----

def extract_structured_recommendations(questions, fragment):
    output = []
    for q in questions:
        af = q["additional_fields"]
        recs = af.get("recommendations")
        answer_options = af.get("answer_options", [])
        tools = af.get("tools", [])

        if fragment in ["B", "C"]:
            if recs:
                output.append({**q, "recommendations": recs})
        elif fragment == "A":
            if recs:
                converted = prompt_convert_fragment_A(q)
                output.append({**q, "recommendations": converted})
        elif fragment in ["D", "E", "F"]:
            generated = prompt_generate_recommendation(q)
            output.append({**q, "recommendations": generated})
    return output

# ---- GRAPH GENERATION ----

def build_fragment_steps(questions):
    steps = []
    for q in questions:
        recs = q.get("recommendations", [])
        if isinstance(recs, str):
            continue  # skip unstructured responses

        for i, step in enumerate(recs):
            steps.append({
                "phase": f"Phase {i+1}",
                "question": q["question"],
                "step_order": i + 1,
                "step_desc": step.get("description", ""),
                "hrs": int(step.get("man_hour_total", 0))
            })
    return steps

def plot_fragment_graph(steps, fragment):
    start = datetime.today()
    timeline = []
    current_start = start

    for step in steps:
        duration = max(1, step['hrs'] // 8)
        end = current_start + timedelta(days=duration)
        timeline.append({
            "Task": f"{step['phase']}: {step['step_desc'][:50]}",
            "Start": current_start,
            "Finish": end,
            "Phase": step['phase']
        })
        current_start = end + timedelta(days=1)

    fig = px.timeline(timeline, x_start="Start", x_end="Finish", y="Task", color="Phase")
    fig.update_yaxes(autorange="reversed")
    fig.update_layout(title=f"Roadmap - Fragment {fragment}", plot_bgcolor='white')
    return fig

# ---- POWERPOINT GENERATION ----

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]

def split_text_into_chunks(text, max_chars=1800):
    return textwrap.wrap(text, max_chars, break_long_words=False, replace_whitespace=False)

def add_text_slide(title, content, font_size=Pt(14)):
    chunks = split_text_into_chunks(content)
    for i, chunk in enumerate(chunks):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        if slide.shapes.title:
            slide.shapes.title.text = f"{title} (cont.)" if i > 0 else title

        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        p = text_frame.paragraphs[0]
        p.text = chunk
        p.font.size = font_size
        for p in text_frame.paragraphs[1:]:
            p.font.size = font_size

def add_graph_slide(title, fig):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    title_shape = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True

    # Save to file then load as image
    temp_path = "temp_graph.png"
    fig.write_image(temp_path)
    with open(temp_path, "rb") as f:
        image_stream = io.BytesIO(f.read())

    pic_left = Inches(1)
    pic_top = Inches(1.2)
    pic_width = Inches(8)
    slide.shapes.add_picture(image_stream, pic_left, pic_top, width=pic_width)

    # Clean up temp file
    if os.path.exists(temp_path):
        os.remove(temp_path)

def add_title_slide(title):
    slide = prs.slides.add_slide(title_slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Fragment Based Roadmap"

# ---- MAIN LOOP ----

for fragment, questions in fragment_groups.items():
    structured_qs = extract_structured_recommendations(questions, fragment)
    steps = build_fragment_steps(structured_qs)

    add_title_slide(f"Fragment {fragment}")

    if steps and fragment in ["B", "C", "A"]:
        summary_text = prompt_summarize_fragment(fragment, structured_qs)
        add_text_slide(f"Roadmap Summary - Fragment {fragment}", summary_text)

        graph = plot_fragment_graph(steps, fragment)
        add_graph_slide(f"Roadmap Chart - Fragment {fragment}", graph)
    else:
        for q in structured_qs:
            recs = q['recommendations']
            content = f"Q: {q['question']}\n\nAnswer: {q['answer']}\n\nRecommendations:\n{recs if isinstance(recs, str) else str(recs)}"
            add_text_slide(f"Fragment {fragment} - Individual Rec", content)

# ---- EXPORT ----

pptx_output = "fragment_roadmap_output.pptx"
prs.save(pptx_output)
print(f"Presentation saved to {pptx_output}")
