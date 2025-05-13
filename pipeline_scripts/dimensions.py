import os
import traceback # Import traceback for detailed error printing
from pptx import Presentation
from pptx.util import Inches, Pt, Cm # Import Inches or Cm

# --- Configuration ---
TEMPLATE_FILENAME = "template.pptx" # Assumes it's in the same directory

# --- !!! IMPORTANT: Specify the Slide Number to Inspect !!! ---
# Remember: Slide numbers in PowerPoint are 1-based (Slide 1, Slide 2, etc.)
# Python list indices are 0-based (Index 0, Index 1, etc.)
# So, to inspect Slide 1, use SLIDE_NUMBER_TO_INSPECT = 1
SLIDE_NUMBER_TO_INSPECT = 2 # <<< CHANGE THIS SLIDE NUMBER AS NEEDED
# ---

def get_slide_placeholder_dimensions(template_path: str, slide_number: int):
    """
    Loads a presentation template, accesses a specific slide by its number,
    and prints the dimensions of all placeholders on that slide.
    """
    if not os.path.exists(template_path):
        print(f"Error: Template file not found at '{template_path}'")
        return

    # Convert 1-based slide number to 0-based index
    slide_index = slide_number - 1

    try:
        prs = Presentation(template_path)
        print(f"Loaded template: '{template_path}'")

        if slide_index < 0 or slide_index >= len(prs.slides):
            print(f"Error: Slide number {slide_number} (index {slide_index}) is out of range.")
            print(f"Presentation has {len(prs.slides)} slides (indices 0 to {len(prs.slides) - 1}).")
            return

        slide = prs.slides[slide_index]
        print(f"\n--- Inspecting Slide Number: {slide_number} (Index: {slide_index}, Layout: '{slide.slide_layout.name}') ---")

        # Check shapes on the slide, filtering for placeholders
        placeholders_on_slide = [shp for shp in slide.shapes if shp.is_placeholder]

        if not placeholders_on_slide:
            print("  No placeholders found on this specific slide.")
            # Optional: List all shapes if needed for debugging
            # print(f"  Total shapes on slide: {len(slide.shapes)}")
            # for i, shp in enumerate(slide.shapes):
            #     print(f"    Shape {i}: Type={shp.shape_type}, Name='{shp.name}'")
            return

        print(f"  Found {len(placeholders_on_slide)} placeholder(s):")
        # Note: The index printed here (i) is just the order found on the slide,
        # it might not match the placeholder index from the layout master.
        for i, ph in enumerate(placeholders_on_slide):
            try:
                # Use .inches directly for cleaner code
                left_in = ph.left.inches
                top_in = ph.top.inches
                width_in = ph.width.inches
                height_in = ph.height.inches

                print(f"\n  Placeholder Found (Order {i}):")
                # Placeholder names might be empty, handle gracefully
                print(f"    Name: '{ph.name if ph.name else '(No Name)'}'")
                # Attempt to print type and original index from layout if available
                try:
                    print(f"    Type: {ph.placeholder_format.type}")
                    print(f"    Layout Index (idx): {ph.placeholder_format.idx}") # Original index from layout
                except AttributeError:
                    print("    Type/Idx: (Could not determine)")

                print(f"    Position & Size (Inches):")
                print(f"      Left:   {left_in:.2f} in")
                print(f"      Top:    {top_in:.2f} in")
                print(f"      Width:  {width_in:.2f} in")
                print(f"      Height: {height_in:.2f} in")
                print("-" * 20)

            except Exception as e_ph:
                print(f"\n  Error processing placeholder {i} on slide {slide_number}: {e_ph}")

    except Exception as e:
        print(f"\nAn error occurred: {e}")
        traceback.print_exc()

# --- Run the function ---
if __name__ == "__main__":
    # Construct the full path assuming the template is in the same directory
    try:
        script_dir = os.path.dirname(os.path.abspath(__file__))
        template_full_path = os.path.join(script_dir, TEMPLATE_FILENAME)
        get_slide_placeholder_dimensions(template_full_path, SLIDE_NUMBER_TO_INSPECT)
    except NameError:
         print("Error: Could not determine script directory (__file__ not defined).")
         print(f"Please ensure '{TEMPLATE_FILENAME}' is in the current working directory.")
         # Try running from current working directory as a fallback
         if os.path.exists(TEMPLATE_FILENAME):
             get_slide_placeholder_dimensions(TEMPLATE_FILENAME, SLIDE_NUMBER_TO_INSPECT)
         else:
             print(f"Error: '{TEMPLATE_FILENAME}' not found in current directory either.")

