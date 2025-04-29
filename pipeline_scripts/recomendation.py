import os
import json
import time
import re
from collections import defaultdict, deque
import traceback
from typing import List, Dict, Any, Optional, Tuple
import sys
import io # To save matplotlib plot to memory and for PPT buffer
import argparse

# --- Third-party Libraries ---
import google.generativeai as genai
import google.api_core.exceptions
from supabase import create_client, Client
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN # Import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt


# --- Configuration Loading ---
load_dotenv()
try:
    current_dir = os.path.abspath(os.path.dirname(__file__))
    project_root = os.path.abspath(os.path.join(current_dir, ".."))
    if project_root not in sys.path: sys.path.append(project_root)
    from backend.utils import config
    SUPABASE_URL = config.SUPABASE_URL
    SUPABASE_KEY = config.SUPABASE_KEY
    GOOGLE_API_KEY = config.GOOGLE_API_KEY
    print("Config loaded from backend.utils")
except (ImportError, ModuleNotFoundError, NameError, AttributeError) as e:
    print(f"Warning: Could not load config from backend.utils ({e}). Falling back to environment variables.")
    SUPABASE_URL = os.getenv("SUPABASE_URL")
    SUPABASE_KEY = os.getenv("SUPABASE_KEY")
    GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
if not all([SUPABASE_URL, SUPABASE_KEY, GOOGLE_API_KEY]):
    print("Error: Missing required environment variables (SUPABASE_URL, SUPABASE_KEY, GOOGLE_API_KEY). Exiting.")
    sys.exit(1)

# --- Global Clients and Settings ---
try: supabase: Client = create_client(SUPABASE_URL, SUPABASE_KEY); print("Supabase client initialized.")
except Exception as e: print(f"Error initializing Supabase client: {e}"); sys.exit(1)
try: genai.configure(api_key=GOOGLE_API_KEY); print("Google GenAI configured.")
except Exception as e: print(f"Error configuring Google GenAI: {e}"); sys.exit(1)

LLM_MODEL_NAME = "gemini-1.5-flash-latest"

# --- Rate Limiting Globals ---
LLM_CALL_TIMESTAMPS = deque(); LLM_TOKEN_TIMESTAMPS = deque()
MAX_RPM = 30; TARGET_RPM = MAX_RPM - 3
TPM_LIMIT = 15000; TARGET_TPM = TPM_LIMIT * 0.9
RPM_WINDOW_SECONDS = 60; MAX_RETRIES = 5; INITIAL_BACKOFF_SECONDS = 2

# --- Constants ---
MAX_CHUNK_TOKENS = 7500; HOURS_PER_DAY = 8
DEFAULT_LAYOUT_TITLE_CONTENT = 1; DEFAULT_LAYOUT_SECTION_HEADER = 2
DEFAULT_LAYOUT_TITLE_ONLY = 5; DEFAULT_LAYOUT_BLANK = 6


# --- LLM Call Function (Modified for TPM) ---
def call_llm(prompt: str, max_tokens: Optional[int] = None, temperature: float = 0.4) -> Optional[str]:
    """Calls the configured Gemini model with rate limiting (RPM & TPM) and retry logic."""
    global LLM_CALL_TIMESTAMPS, LLM_TOKEN_TIMESTAMPS
    current_time = time.monotonic()
    # RPM Check
    while LLM_CALL_TIMESTAMPS and LLM_CALL_TIMESTAMPS[0] <= current_time - RPM_WINDOW_SECONDS: LLM_CALL_TIMESTAMPS.popleft()
    if len(LLM_CALL_TIMESTAMPS) >= TARGET_RPM:
        wait_time = max(0, (LLM_CALL_TIMESTAMPS[0] + RPM_WINDOW_SECONDS) - current_time)
        print(f"RPM limit ({len(LLM_CALL_TIMESTAMPS)}/{TARGET_RPM}). Wait {wait_time:.2f}s..."); time.sleep(wait_time + 0.1); current_time = time.monotonic()
        while LLM_CALL_TIMESTAMPS and LLM_CALL_TIMESTAMPS[0] <= current_time - RPM_WINDOW_SECONDS: LLM_CALL_TIMESTAMPS.popleft()
    # TPM Check
    estimated_output_tokens = max_tokens if max_tokens is not None else 2048; estimated_input_tokens = 0
    try: model_for_counting = genai.GenerativeModel(LLM_MODEL_NAME); estimated_input_tokens = model_for_counting.count_tokens(prompt).total_tokens
    except Exception as e: print(f"Warning: Token count failed ({e}). Fallback."); estimated_input_tokens = len(prompt) // 4
    current_call_estimated_tokens = estimated_input_tokens + estimated_output_tokens
    while LLM_TOKEN_TIMESTAMPS and LLM_TOKEN_TIMESTAMPS[0][0] <= current_time - RPM_WINDOW_SECONDS: LLM_TOKEN_TIMESTAMPS.popleft()
    tokens_in_window = sum(t[1] for t in LLM_TOKEN_TIMESTAMPS)
    while tokens_in_window + current_call_estimated_tokens > TARGET_TPM and LLM_TOKEN_TIMESTAMPS:
        wait_time = max(0, (LLM_TOKEN_TIMESTAMPS[0][0] + RPM_WINDOW_SECONDS) - current_time)
        print(f"TPM limit ({(tokens_in_window + current_call_estimated_tokens):.0f}/{TARGET_TPM:.0f} est). Wait {wait_time:.2f}s..."); time.sleep(wait_time + 0.1); current_time = time.monotonic()
        while LLM_TOKEN_TIMESTAMPS and LLM_TOKEN_TIMESTAMPS[0][0] <= current_time - RPM_WINDOW_SECONDS: LLM_TOKEN_TIMESTAMPS.popleft()
        tokens_in_window = sum(t[1] for t in LLM_TOKEN_TIMESTAMPS)
    # API Call
    intended_start_time = current_time; LLM_CALL_TIMESTAMPS.append(intended_start_time)
    backoff_time = INITIAL_BACKOFF_SECONDS
    for attempt in range(MAX_RETRIES):
        try:
            print(f"Attempt {attempt + 1}/{MAX_RETRIES}: Calling LLM (Input est: ~{estimated_input_tokens} tokens)...")
            llm = genai.GenerativeModel(LLM_MODEL_NAME)
            gen_config = genai.types.GenerationConfig(temperature=temperature, max_output_tokens=max_tokens if max_tokens is not None else None)
            safety_settings=[{"category": c, "threshold": "BLOCK_MEDIUM_AND_ABOVE"} for c in ["HARM_CATEGORY_HARASSMENT", "HARM_CATEGORY_HATE_SPEECH", "HARM_CATEGORY_SEXUALLY_EXPLICIT", "HARM_CATEGORY_DANGEROUS_CONTENT"]]
            response = llm.generate_content(prompt, generation_config=gen_config, safety_settings=safety_settings)
            if response.parts:
                total_tokens_used = 0
                try: total_tokens_used = response.usage_metadata.total_token_count; print(f"LLM Success. Tokens: In={response.usage_metadata.prompt_token_count}, Out={response.usage_metadata.candidates_token_count}, Total={total_tokens_used}")
                except (AttributeError, TypeError, ValueError): total_tokens_used = estimated_input_tokens + (len(response.text) // 4); print(f"LLM Success. Estimated Tokens: Total={total_tokens_used:.0f}")
                LLM_TOKEN_TIMESTAMPS.append((time.monotonic(), total_tokens_used)); return response.text.strip()
            elif response.prompt_feedback.block_reason: print(f"Warning: Blocked. Reason: {response.prompt_feedback.block_reason}"); break
            else: print("Warning: LLM returned no content or block reason."); break
        except google.api_core.exceptions.ResourceExhausted as e: print(f"Error: ResourceExhausted attempt {attempt + 1}: {e}");
        except Exception as e: print(f"Error during LLM call attempt {attempt + 1}: {e}"); traceback.print_exc(); break
        if LLM_CALL_TIMESTAMPS and LLM_CALL_TIMESTAMPS[-1] == intended_start_time: LLM_CALL_TIMESTAMPS.pop()
        if attempt < MAX_RETRIES - 1: print(f"Waiting {backoff_time:.2f}s before retry..."); time.sleep(backoff_time); backoff_time *= 2; intended_start_time = time.monotonic(); LLM_CALL_TIMESTAMPS.append(intended_start_time)
        else: print("Max retries reached."); return None
    if LLM_CALL_TIMESTAMPS and LLM_CALL_TIMESTAMPS[-1] == intended_start_time: LLM_CALL_TIMESTAMPS.pop()
    return None

# --- Database Interaction Functions ---
def get_supabase_client() -> Client:
    global supabase;
    if not supabase: supabase = create_client(SUPABASE_URL, SUPABASE_KEY)
    return supabase
def fetch_followup_data(supabase_client: Client, session_id: str) -> List[Dict[str, Any]]:
    try: response = supabase_client.table("followup_responses").select("*").eq("session_id", session_id).execute(); return response.data or []
    except Exception as e: print(f"Error fetching followup: {e}"); traceback.print_exc(); return []
def fetch_question_details(supabase_client: Client, question_ids: List[int]) -> Dict[int, Dict[str, Any]]:
    if not question_ids: return {}
    try: valid_question_ids = [int(qid) for qid in question_ids if isinstance(qid, (int, str)) and str(qid).isdigit()]; response = supabase_client.table("questions").select("id, category, subcategory, sr_no, fragments, question, additional_fields, text_representation").in_("id", valid_question_ids).execute(); return {q['id']: q for q in response.data} if response.data else {}
    except Exception as e: print(f"Error fetching questions: {e}"); traceback.print_exc(); return {}

# MODIFIED: Fetch frag_names as well
def fetch_fragment_info(supabase_client: Client, fragment_ids: List[str]) -> Dict[str, Dict[str, Any]]:
    """Fetches fragment descriptions and names for a list of CLEANED fragment IDs."""
    if not fragment_ids: return {}
    try:
        unique_fragment_ids = list(set(filter(None, fragment_ids)))
        if not unique_fragment_ids: print("Warning: No valid unique fragment IDs provided."); return {}
        # Select fragment, description, category, AND frag_names
        response = supabase_client.table("fragment_info").select(
            "fragment, description, category, frag_names"
        ).in_("fragment", unique_fragment_ids).execute()
        if response.data:
            # Store fragment_id -> {all_data}
            return {f['fragment']: f for f in response.data}
        else:
            print(f"Warning: No fragment info found matching fragments: {unique_fragment_ids}.")
            # Return default structure even if not found, so grouping doesn't fail
            return {frag_id: {"fragment": frag_id, "description": "N/A", "category": "N/A", "frag_names": f"Fragment {frag_id}"} for frag_id in unique_fragment_ids}
    except Exception as e:
        print(f"Error fetching fragment info: {e}"); traceback.print_exc()
        # Return default structure on error
        return {frag_id: {"fragment": frag_id, "description": "N/A", "category": "N/A", "frag_names": f"Fragment {frag_id}"} for frag_id in fragment_ids}


# --- Data Processing Functions ---
def process_additional_fields(question: Dict[str, Any], answer_data: Dict[str, Any]) -> Dict[str, Any]:
    # ... (no changes needed in this function) ...
    additional_fields = question.get("additional_fields");
    if not additional_fields: return {}
    if isinstance(additional_fields, str):
        try: additional_fields = json.loads(additional_fields)
        except json.JSONDecodeError: return {}
    if not isinstance(additional_fields, dict): return {}
    filtered_fields = {}; user_answer = answer_data.get("value") if isinstance(answer_data, dict) else None
    always_include_keys = ["guidelines", "answer_options", "sub_sub_category", "subjective_answer", "multiple_correct_answer_options", "why_use_these_tools"]
    for key in always_include_keys:
        if key in additional_fields: filtered_fields[key] = additional_fields[key]
    retrieve_all_fields = False; processed_specifically = False
    if user_answer == "Implemented": processed_specifically = True; filtered_fields.pop("recommendation", None); filtered_fields.pop("recommendations", None); filtered_fields.pop("tools", None)
    elif user_answer == "Not Implemented": retrieve_all_fields = True; processed_specifically = True
    elif user_answer == "dont know" or "subjective_answer" in additional_fields or "multiple_correct_answer_options" in additional_fields: retrieve_all_fields = True; processed_specifically = True
    elif "recommendations" in additional_fields and isinstance(additional_fields["recommendations"], dict):
        matched_recommendation = None; match_found_in_b = False; normalized_answer = user_answer.strip().lower() if user_answer else ""
        for rec_key, rec_value in additional_fields["recommendations"].items():
            normalized_key = rec_key.strip().lower(); key_base = re.sub(r'\s*\(level\s*\d+\)$', '', normalized_key).strip()
            if normalized_answer == normalized_key or (key_base and normalized_answer == key_base) or (key_base and normalized_answer.startswith(key_base)): matched_recommendation = {rec_key: rec_value}; match_found_in_b = True; break
        if match_found_in_b: processed_specifically = True; filtered_fields["recommendations"] = matched_recommendation;
        if "tools" in additional_fields: filtered_fields["tools"] = additional_fields["tools"]
        else: retrieve_all_fields = True
    if not processed_specifically: retrieve_all_fields = True
    if retrieve_all_fields: temp_copy = additional_fields.copy(); temp_copy.update(filtered_fields); filtered_fields = temp_copy
    return {k: v for k, v in filtered_fields.items() if v is not None}

# MODIFIED: To use the fetched fragment_info containing frag_names
def group_data_by_fragment(followup_responses: List[Dict[str, Any]], questions_details: Dict[int, Dict[str, Any]]) -> Dict[str, Dict[str, Any]]:
    """Groups processed question data by fragment, ensuring fragment_info is attached."""
    grouped_data = defaultdict(lambda: {"fragment_info": None, "questions_data": []});
    fragment_ids_found = set()

    # First pass: Group questions by fragment ID
    for resp in followup_responses:
        q_id = resp.get("question_id"); question_detail = questions_details.get(q_id);
        if not question_detail: continue
        fragment_list = question_detail.get("fragments"); raw_fragment_id = None; cleaned_fragment_id = None
        if isinstance(fragment_list, (list, tuple)) and fragment_list: raw_fragment_id = str(fragment_list[0]).strip()
        elif isinstance(fragment_list, str): raw_fragment_id = fragment_list.strip()
        else: continue
        if raw_fragment_id: match = re.match(r'^[\[\'\"]*(.*?)[\'\"\]]*$', raw_fragment_id); cleaned_fragment_id = match.group(1) if match else raw_fragment_id
        if not cleaned_fragment_id: continue

        fragment_ids_found.add(cleaned_fragment_id) # Collect all unique fragment IDs

        answer_json = resp.get("answer")
        if isinstance(answer_json, str):
            try: answer_json = json.loads(answer_json)
            except json.JSONDecodeError: answer_json = {}
        elif not isinstance(answer_json, dict): answer_json = {}
        processed_add_fields = process_additional_fields(question_detail, answer_json)
        question_data = {"question_id": q_id, "question_text": question_detail.get("question"), "category": question_detail.get("category"), "subcategory": question_detail.get("subcategory"), "user_answer": answer_json, "processed_additional_fields": processed_add_fields,}
        grouped_data[cleaned_fragment_id]["questions_data"].append(question_data)

    # Second pass: Fetch info for all found fragments and attach it
    if fragment_ids_found:
        all_fragment_info = fetch_fragment_info(get_supabase_client(), list(fragment_ids_found))
        for frag_id, data in grouped_data.items():
            # Use fetched info, provide default if lookup fails for some reason
            data["fragment_info"] = all_fragment_info.get(
                frag_id,
                {"fragment": frag_id, "description": "N/A", "category": "N/A", "frag_names": f"Fragment {frag_id}"}
            )
    else: print("No valid fragments found to fetch info for.")

    return dict(grouped_data)

# --- LLM Prompt Generation Functions ---
# MODIFIED: To use frag_names from fragment_info
def format_chunk_for_prompt(fragment_info: Dict[str, Any], questions_chunk: List[Dict[str, Any]]) -> str:
    """Formats a chunk of questions for a fragment into a string for LLM prompts."""
    output = []
    # Use frag_names if available, otherwise fallback
    frag_display_name = fragment_info.get('frag_names') or fragment_info.get('category') or fragment_info.get('fragment', 'N/A')
    output.append(f"Fragment Name: {frag_display_name}") # Use Name instead of ID
    output.append(f"Fragment Description: {fragment_info.get('description', 'N/A')}")
    output.append(f"(Internal Fragment ID: {fragment_info.get('fragment', 'N/A')})") # Keep ID for internal reference if needed by LLM? Optional.
    output.append("\n--- User Assessment Data (Chunk) ---")
    for q_data in questions_chunk:
        output.append(f"\nQuestion ID: {q_data['question_id']} (Topic: {q_data.get('subcategory', 'N/A')})")
        output.append(f"Question: {q_data['question_text']}")
        try: user_answer_str = json.dumps(q_data['user_answer'])
        except TypeError: user_answer_str = str(q_data['user_answer'])
        output.append(f"User Answer: {user_answer_str}")
        proc_fields = q_data.get('processed_additional_fields', {})
        if proc_fields:
            output.append("  Relevant Context/Metadata:")
            if 'guidelines' in proc_fields: output.append(f"    - Guideline: {proc_fields['guidelines']}")
            if 'sub_sub_category' in proc_fields: output.append(f"    - Sub-Topic: {proc_fields['sub_sub_category']}")
            try:
                if 'tools' in proc_fields: output.append(f"    - Available Tools Info: {json.dumps(proc_fields['tools'], indent=2)}")
                if 'recommendations' in proc_fields: output.append(f"    - Specific Recommendation Context: {json.dumps(proc_fields['recommendations'], indent=2)}")
                elif 'recommendation' in proc_fields: output.append(f"    - General Recommendation Context: {json.dumps(proc_fields['recommendation'], indent=2)}")
            except (TypeError, ValueError) as json_err: output.append(f"    - (Error formatting context field: {json_err})")
        output.append("-" * 10)
    return "\n".join(output)

def generate_insights_prompt(fragment_context: str) -> str:
    """Creates the prompt for generating Top 3 Insights (used for chunks). Refined."""
    # ... (prompt content same as previous version - already avoids IDs) ...
    prompt = f"""
Based on the following assessment data for a specific fragment chunk, identify up to 3 key insights regarding the user's current state, challenges, or strengths related to the topics covered in this chunk.

Assessment Data Chunk:
{fragment_context}

Instructions:
- Analyze the questions, user answers, and associated metadata within this chunk.
- Synthesize the information to derive meaningful insights specific to this chunk.
- **Insights must be specific, actionable, and avoid vague jargon.**
- **Do NOT refer to specific Question IDs.** Instead, refer to the themes or topics assessed (e.g., "Regarding version control practices...", "The assessment of monitoring tools indicates...").
- Present the insights as distinct points or paragraphs (up to 3).
- Focus on the most impactful observations based *only* on the provided chunk data.
- Ensure the output is formatted clearly, starting each insight on a new line.
- Do not include introductory phrases like "Here are the insights:".
- Do not use markdown formatting like '***', '---', or numbering like '1.'.

Top Insights for this Chunk:
[Insight 1 text]
[Insight 2 text]
[Insight 3 text]
"""
    return prompt

def generate_roadmap_prompt(fragment_context: str) -> str:
    """Creates the prompt for generating the Roadmap Summary as JSON (used for chunks). Refined."""
    # ... (prompt content same as previous version - already avoids IDs) ...
    prompt = f"""
Based *only* on the following assessment data chunk for a specific fragment, generate a 3-phase implementation roadmap (Early, Intermediate, Advanced Steps) **strictly as a JSON object**.

Assessment Data Chunk:
{fragment_context}

Instructions:
- Analyze the questions, user answers, and available metadata *within this chunk*.
- Synthesize this information to create a coherent *partial* roadmap based *only* on this chunk's data.
- **Aggregate existing step-wise recommendations within this chunk**: Merge Step 1s into 'Early Steps', Step 2s into 'Intermediate Steps', etc.
- **Format `how_methodology`:** For each phase, create a single string containing a sequentially numbered list of specific, actionable steps (e.g., "1. Implement basic monitoring for critical services. 2. Define initial SLOs based on current performance. 3. Setup automated alerts for SLO breaches."). **Do NOT repeat 'Step X:' labels for each item.**
- **Estimate man-hours realistically** based on data *in this chunk*. Ensure 'total_hours' per phase is a numeric integer.
- **Consolidate roles and tools** mentioned *in this chunk* for each phase.
- **Content Specificity:** Ensure the `summary_context`, phase `description`s, and `how_methodology` steps are specific and avoid vague jargon.
- **No Question IDs:** Do NOT refer to specific Question IDs in any field. Refer to topics/themes if necessary.
- **Output ONLY the JSON object**, starting with `{{` and ending with `}}`. No extra text or markdown.

Required JSON Structure (for this partial roadmap):
{{
  "fragment_id": "[Internal Fragment ID from Assessment Data]",
  "fragment_name": "[Fragment Name from Assessment Data]",
  "summary_context": "[Specific ~1-2 sentence narrative overview *based on this chunk*]",
  "total_estimated_hours": "[Numeric integer sum of hours *from this chunk*]",
  "key_roles_involved": "[Consolidated list/string of unique roles *from this chunk*]",
  "phases": [
    {{
      "phase_name": "Early Steps",
      "step_title": "[Specific title, e.g., Foundational Monitoring Setup]",
      "description": "[Specific overview based on this chunk and Step 1s]",
      "how_methodology": "[Single string with numbered list of actions, e.g., '1. Action A. 2. Action B.']",
      "who_roles_involved": "[List/string of roles *from this chunk*]",
      "tools_platforms_used": "[List/string of tools *from this chunk*]",
      "estimate_per_subtask": "[Labelled subtasks/hours string *from this chunk*]",
      "total_hours": "[Numeric integer total hours for this phase *from this chunk*]"
    }},
    // ... Intermediate and Advanced Steps structure same as before ...
     {{
      "phase_name": "Intermediate Steps",
      "step_title": "[Specific title, e.g., Process Integration & Automation]",
      "description": "[Specific goal based on this chunk and Step 2s]",
      "how_methodology": "[Single string with numbered list of actions, e.g., '1. Action C. 2. Action D.']",
      "who_roles_involved": "[List/string of roles *from this chunk*]",
      "tools_platforms_used": "[List/string of tools *from this chunk*]",
      "estimate_per_subtask": "[Labelled subtasks/hours string *from this chunk*]",
      "total_hours": "[Numeric integer total hours for this phase *from this chunk*]"
    }},
    {{
      "phase_name": "Advanced Steps",
      "step_title": "[Specific title, e.g., Performance Optimization & Governance]",
      "description": "[Specific plan based on this chunk and Step 3s]",
      "how_methodology": "[Single string with numbered list of actions, e.g., '1. Action E. 2. Action F.']",
      "who_roles_involved": "[List/string of roles *from this chunk*]",
      "tools_platforms_used": "[List/string of tools *from this chunk*]",
      "estimate_per_subtask": "[Labelled subtasks/hours string *from this chunk*]",
      "total_hours": "[Numeric integer total hours for this phase *from this chunk*]"
    }}
  ]
}}
"""
    return prompt

def generate_synthesize_insights_prompt(partial_insights_list: List[str]) -> str:
    """Creates the prompt for synthesizing partial insights into the final Top 3. Refined."""
    # ... (prompt content same as previous version - already avoids IDs) ...
    combined_insights_text = "\n\n---\n\n".join(f"Partial Insights Set {i+1}:\n{insight}" for i, insight in enumerate(partial_insights_list) if insight and insight.strip())
    if not combined_insights_text: return ""
    prompt = f"""
Review the following sets of partial insights generated from analyzing different chunks of assessment data for the *same* overall fragment. Your task is to synthesize these partial observations into the definitive Top 3 key insights for the entire fragment.

Combined Partial Insights:
{combined_insights_text}

Instructions:
- Identify recurring themes, major challenges, and significant strengths across all partial insights.
- Consolidate overlapping points.
- Prioritize the most impactful and overarching observations for the final Top 3.
- **Final insights must be specific, actionable, and avoid vague jargon.**
- **Do NOT refer to specific Question IDs.** Instead, refer to the themes or topics assessed (e.g., "Overall, the approach to automated testing lacks...", "A key strength identified is the existing CI/CD pipeline...").
- Present the final insights as 3 distinct points or paragraphs.
- Focus on clarity and conciseness.
- Do not include introductory phrases like "Here are the final insights:". Just provide the insights directly.
- Do not use markdown formatting like '***', '---', or numbering like '1.'. Start each final insight on a new line.

Final Top 3 Insights:
[Final Insight 1 text]
[Final Insight 2 text]
[Final Insight 3 text]
"""
    return prompt

# MODIFIED: Synthesis Prompt for Roadmap (V4 - Paragraph Format, Summarization, No summary_context, No who_roles)
def generate_synthesize_roadmaps_prompt(partial_roadmaps_list: List[Dict[str, Any]], fragment_id: str, fragment_name: str) -> str:
    """
    Creates the prompt for synthesizing partial roadmaps into a FINAL,
    CONSOLIDATED roadmap JSON structure. Removes summary_context and consolidated who_roles.
    """
    valid_partial_roadmaps = [r for r in partial_roadmaps_list if r and isinstance(r, dict) and r.get("phases")]
    if not valid_partial_roadmaps: return ""
    try: partial_roadmaps_json_str = json.dumps(valid_partial_roadmaps, indent=2)
    except TypeError as e: print(f"Error serializing partial roadmaps: {e}"); 
    try: partial_roadmaps_json_str = json.dumps([str(r) for r in valid_partial_roadmaps], indent=2); 
    except Exception as e2: print(f"Fallback serialization failed: {e2}"); return ""

    prompt = f"""
You are given a list of partial implementation roadmaps (as a JSON array), each generated from analyzing a different chunk of assessment data for the *same* overall fragment (ID: {fragment_id}, Name: {fragment_name}). Your task is to synthesize these into a single, **consolidated roadmap description** while **preserving the distinct hour totals** for the original Early, Intermediate, and Advanced phases for graphing purposes. Output **strictly as a JSON object** following the specified structure.

List of Partial Roadmaps (JSON Array):
{partial_roadmaps_json_str}

Instructions for Synthesis and Final JSON Output:
1.  **Analyze All Partial Roadmaps:** Review the content within the 'Early Steps', 'Intermediate Steps', and 'Advanced Steps' phases of each partial roadmap.
2.  **Calculate Phase Hour Totals:** Sum the `total_hours` for all 'Early Steps' phases across all partial roadmaps. Do the same for 'Intermediate Steps' and 'Advanced Steps'. Store these three sums accurately as numeric integers.
3.  **Consolidate Content:**
    * **Synthesize `how_methodology`:** Critically review the `how_methodology` steps from all partial roadmaps for the corresponding phase. **Merge these into a single string formatted as a paragraph.** Inject clear inline headers (e.g., "\\n--- Early Steps ---\\n", "\\n--- Intermediate Steps ---\\n", "\\n--- Advanced Steps ---\\n") before each phase's steps. **Within each phase's section, format the steps as a single flowing paragraph with sequential numbering (e.g., '1. Action A. 2. Action B. 3. Action C.') using spaces after the period, NOT newlines.** Where appropriate and without losing critical information, **combine very short, directly related steps into a single numbered item** for conciseness (e.g., '4. Setup basic monitoring and alerting.'). Ensure steps remain specific and actionable. **Do NOT repeat 'Step X:' labels.**
    * **Synthesize `estimate_per_subtask`:** Similarly, combine the subtask estimates from all partial phases into a **single string formatted as a paragraph.** Group estimates under inline phase headers ("--- Early Steps ---", etc.). **Separate individual estimates within each phase using semicolons and spaces (e.g., 'Task A - 10h; Task B - 15h; Task C - 5h'). Avoid using newlines for each estimate.** Aggregate hours for similar tasks where appropriate. Ensure clear labeling.
    * **Consolidate `tools_platforms_used`:** Create a single list/string of unique tools mentioned across *all* phases in *all* partial roadmaps.
    * **Synthesize `description`:** Write a single, coherent narrative description that flows through the objectives and activities originally planned for the early, intermediate, and advanced stages. Aim for specificity and reasonable conciseness. Avoid jargon.
    * **Create Consolidated `title`:** Create a single, overarching title for the entire consolidated roadmap (e.g., "{fragment_name} Implementation Plan").
4.  **Generate Top-Level Fields:**
    * **`fragment_id` & `fragment_name`:** Use the provided values: "{fragment_id}", "{fragment_name}".
    * **`total_estimated_hours`:** Sum the three calculated phase hour totals from step 2. Ensure numeric integer.
    * **`key_roles_involved`:** Create a final consolidated list/string of all unique roles mentioned in the `who_roles_involved` or `key_roles_involved` fields across *all* partial roadmaps.
5.  **Structure Output JSON:** Create the final JSON object with the exact structure specified below. Place the calculated phase hour totals (as numeric integers) within the `graph_phase_data` field. Place all the consolidated/merged content within the `consolidated_details` field. **Do NOT include `summary_context` or `who_roles_involved` inside `consolidated_details`.**
6.  **Content Quality:** Ensure all text is specific, actionable, and avoids vague language or excessive length. **Do NOT refer to specific Question IDs.**
7.  **Output ONLY JSON:** Output *only* the final JSON object, starting with `{{` and ending with `}}`. No extra text, markdown, or explanations.

Required Final JSON Structure:
{{
  "fragment_id": "{fragment_id}",
  "fragment_name": "{fragment_name}",
  "total_estimated_hours": "[Numeric integer sum of all phase hours]",
  "key_roles_involved": "[Final consolidated list/string of unique roles across all phases]",
  "graph_phase_data": {{
    "Early Steps": {{"total_hours": "[Numeric integer sum of hours for ALL Early Steps]"}},
    "Intermediate Steps": {{"total_hours": "[Numeric integer sum of hours for ALL Intermediate Steps]"}},
    "Advanced Steps": {{"total_hours": "[Numeric integer sum of hours for ALL Advanced Steps]"}}
  }},
  "consolidated_details": {{
    "title": "[Consolidated title for the overall roadmap]",
    "description": "[Consolidated narrative description covering progression through phases]",
    "how_methodology": "[SINGLE string PARAGRAPH with inline phase headers and merged, sequentially re-numbered steps separated by SPACES]",
    "tools_platforms_used": "[Consolidated list/string of unique tools for the entire roadmap]",
    "estimate_per_subtask": "[SINGLE string PARAGRAPH with inline phase headers and aggregated/merged subtask estimates separated by SEMICOLONS]"
  }}
}}
"""
    return prompt

# --- Roadmap Parsing (Updated for New Synthesis Structure V2) ---
def parse_roadmap_output_json(roadmap_json_text: str) -> Optional[Dict[str, Any]]:
    """
    Parses the CONSOLIDATED roadmap JSON structure (V2) from the synthesis step.
    Extracts graph data and consolidated text details (excluding summary_context and who_roles).
    """
    if not roadmap_json_text: print("Warning: Roadmap JSON text is empty."); return None
    cleaned_text = roadmap_json_text.strip()
    cleaned_text = re.sub(r'^```(?:json)?\s*', '', cleaned_text, flags=re.IGNORECASE)
    cleaned_text = re.sub(r'\s*```$', '', cleaned_text); cleaned_text = cleaned_text.strip()
    cleaned_text = re.sub(r',\s*([\}\]])', r'\1', cleaned_text)
    try:
        if not cleaned_text.startswith("{") or not cleaned_text.endswith("}"):
            json_match = re.search(r'\{.*\S.*\}', cleaned_text, re.DOTALL)
            if json_match: cleaned_text = json_match.group(0)
            else: print("Error: Could not find valid JSON object structure."); raise json.JSONDecodeError("Invalid JSON structure", cleaned_text, 0)
        roadmap_data = json.loads(cleaned_text)
        if not isinstance(roadmap_data, dict): print("Error: Parsed roadmap JSON is not a dictionary."); return None

        final_data = {}
        # Top-level fields (summary_context removed)
        final_data['fragment_id'] = roadmap_data.get('fragment_id')
        final_data['fragment_name'] = roadmap_data.get('fragment_name')
        # final_data['summary_context'] = roadmap_data.get('summary_context') # REMOVED
        final_data['key_roles_involved'] = roadmap_data.get('key_roles_involved')
        if isinstance(final_data['key_roles_involved'], list): final_data['key_roles_involved'] = ", ".join(map(str, final_data['key_roles_involved']))
        elif final_data['key_roles_involved'] is None: final_data['key_roles_involved'] = "N/A"

        # Graph data extraction
        graph_data = roadmap_data.get('graph_phase_data')
        temp_phases_for_graph = []
        total_hours_from_graph_data = 0
        expected_phases = ["Early Steps", "Intermediate Steps", "Advanced Steps"]
        if isinstance(graph_data, dict):
            for phase_name in expected_phases:
                 phase_info = graph_data.get(phase_name)
                 if isinstance(phase_info, dict) and 'total_hours' in phase_info:
                     try: hours = int(float(phase_info['total_hours'])); temp_phases_for_graph.append({"name": phase_name, "total_hours": hours}); total_hours_from_graph_data += hours
                     except (ValueError, TypeError): print(f"Warning: Invalid hours in graph_phase_data for {phase_name}.")
                 else: print(f"Warning: Missing/invalid phase_info in graph_phase_data for {phase_name}.")
        else: print("Warning: 'graph_phase_data' missing or invalid.")
        final_data['phases_for_graph'] = temp_phases_for_graph

        # Validate top-level total_estimated_hours
        try: final_data['total_estimated_hours'] = int(float(roadmap_data.get('total_estimated_hours', 0)))
        except (ValueError, TypeError): final_data['total_estimated_hours'] = total_hours_from_graph_data
        if final_data.get('total_estimated_hours', 0) != total_hours_from_graph_data:
            print(f"Warning: Top-level hours ({final_data.get('total_estimated_hours')}) != graph phase sum ({total_hours_from_graph_data}). Using sum.")
            final_data['total_estimated_hours'] = total_hours_from_graph_data

        # Consolidated details (who_roles_involved removed)
        consolidated_details = roadmap_data.get('consolidated_details')
        if isinstance(consolidated_details, dict):
             final_data['details'] = {
                 "title": consolidated_details.get("title", "Consolidated Implementation Plan"),
                 "description": consolidated_details.get("description", "N/A"), # Keep description
                 "how_methodology": consolidated_details.get("how_methodology", "N/A"),
                 # "who_roles_involved": consolidated_details.get("who_roles_involved", "N/A"), # REMOVED
                 "tools_platforms_used": consolidated_details.get("tools_platforms_used", "N/A"),
                 "estimate_per_subtask": consolidated_details.get("estimate_per_subtask", "N/A")
             }
             # Convert lists to strings if necessary
             # if isinstance(final_data['details']['who_roles_involved'], list): final_data['details']['who_roles_involved'] = ", ".join(map(str, final_data['details']['who_roles_involved'])) # REMOVED
             if isinstance(final_data['details']['tools_platforms_used'], list): final_data['details']['tools_platforms_used'] = ", ".join(map(str, final_data['details']['tools_platforms_used']))
        else: print("Warning: 'consolidated_details' missing or invalid."); final_data['details'] = {"error": "Consolidated details missing or invalid."}

        if not final_data.get('phases_for_graph') and final_data.get('details', {}).get('error'):
             print("Error: Parsing failed to extract usable graph data or consolidated details."); return None

        return final_data
    except json.JSONDecodeError as e: print(f"Error: Failed to decode roadmap JSON: {e}"); return None
    except Exception as e: print(f"Error processing parsed roadmap data: {e}"); traceback.print_exc(); return None

# --- Graphing Function ---
def create_roadmap_graph(graph_phase_data: List[Dict[str, Any]]) -> Optional[io.BytesIO]:
    """Generates roadmap bar chart from phase data (name, total_hours)."""
    # ... (no changes needed in this function) ...
    if not graph_phase_data or not isinstance(graph_phase_data, list): return None
    phase_names = [p.get("name", f"P{i+1}") for i, p in enumerate(graph_phase_data)]
    total_hours = [p.get("total_hours", 0) for p in graph_phase_data]
    total_days = [h / HOURS_PER_DAY for h in total_hours]
    if not any(d >= 0 for d in total_days): return None
    try:
        fig, ax = plt.subplots(figsize=(7, 4))
        bars = ax.bar(phase_names, total_days, color=['#4A90E2', '#F5A623', '#50E3C2'][:len(phase_names)])
        ax.set_ylabel(f'Est. Duration (Days)', fontsize=9); ax.set_title(f'Estimated Phase Durations', fontsize=10)
        ax.set_xticks(range(len(phase_names))); ax.set_xticklabels(phase_names, rotation=0, ha='center', fontsize=9)
        ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False); ax.yaxis.grid(True, linestyle='--', alpha=0.6); ax.tick_params(axis='y', labelsize=9)
        for bar in bars: yval = bar.get_height();
        if yval > 0: plt.text(bar.get_x() + bar.get_width()/2.0, yval, f'{yval:.1f}', va='bottom', ha='center', fontsize=8)
        elif yval == 0 : plt.text(bar.get_x() + bar.get_width()/2.0, yval, f'0', va='bottom', ha='center', fontsize=8)
        plt.tight_layout(pad=0.5); img_buffer = io.BytesIO(); plt.savefig(img_buffer, format='png', dpi=120); img_buffer.seek(0); plt.close(fig); return img_buffer
    except Exception as e: print(f"Error generating graph: {e}"); traceback.print_exc(); return None

# --- PPT Generation Functions ---
def add_text_to_shape(shape, text):
    """Helper to add text to a shape, handling None values."""
    # ... (no changes needed) ...
    if shape and hasattr(shape, "text_frame"): tf = shape.text_frame; tf.text = str(text) if text is not None else ""; tf.word_wrap = True
    elif shape and hasattr(shape, "text"): shape.text = str(text) if text is not None else ""

# Use the version that applies justification but allows default auto-fitting
def add_text_slide(prs: Presentation, title_text: str, content_text: str, layout_index: int = DEFAULT_LAYOUT_TITLE_CONTENT):
    """Adds slide with title/content, using default auto-fitting, sets content to Justified."""
    # ... (function body same as before - ID: ppt_add_text_slide_justified) ...
    try: slide_layout = prs.slide_layouts[layout_index]
    except IndexError: print(f"Error: Layout index {layout_index} OOR. Using 0."); slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout); title_shape = None; content_shape = None; title_ph_idx = -1; content_ph_idx = -1
    try: title_shape = slide.shapes.title; title_ph_idx = title_shape.placeholder_format.idx
    except AttributeError: pass
    try:
         if len(slide.placeholders) > 1: content_shape = slide.placeholders[1]; content_ph_idx = 1
    except (KeyError, IndexError, AttributeError): pass
    if title_shape is None or content_shape is None:
        for i, ph in enumerate(slide.placeholders):
            ph_name = ph.name.lower() if ph.name else ""; ph_type = ph.placeholder_format.type
            if title_shape is None and ('title' in ph_name or ph.placeholder_format.idx == 0 or ph.placeholder_format.idx == 10): title_shape = ph; title_ph_idx = i
            elif content_shape is None and ('content' in ph_name or 'body' in ph_name or 'text' in ph_name or ph.placeholder_format.idx == 1 or ph.placeholder_format.idx == 11):
                 if i != title_ph_idx: content_shape = ph; content_ph_idx = i
            if title_shape is not None and content_shape is not None: break
    if title_shape: add_text_to_shape(title_shape, title_text)
    else: print(f"Warning: Could not find title placeholder for '{title_text}' on layout {layout_index}.")
    if content_shape:
        tf = content_shape.text_frame; tf.clear(); add_text_to_shape(tf, content_text) # Add text first
        try: # Apply Justification AFTER adding text
            for paragraph in tf.paragraphs:
                if paragraph.text.strip(): paragraph.alignment = PP_ALIGN.JUSTIFY
        except Exception as e_align: print(f"Warning: Could not apply justification to content on slide '{title_text}': {e_align}")
        tf.vertical_anchor = MSO_ANCHOR.TOP; tf.word_wrap = True
    else:
        print(f"Warning: Could not find content placeholder for '{title_text}' on layout {layout_index}. Adding to notes.")
        try: notes_tf = slide.notes_slide.notes_text_frame; notes_tf.text = f"Title: {title_text}\n\nContent:\n{content_text}"
        except Exception as e_notes: print(f"Error adding content to notes slide: {e_notes}")


# --- Modified build_ppt_report function ---
def build_ppt_report(session_id: str, report_data: Dict[str, Dict[str, Any]], template_path_relative: Optional[str] = "template.pptx") -> Optional[io.BytesIO]:
    """
    Builds the PowerPoint presentation in memory.
    Uses fetched frag_names, consolidated roadmap details, removes fragment title slide,
    and deletes the first slide of the template before saving.
    """
    prs = None
    # Template Loading
    if template_path_relative:
        try:
            script_dir = os.path.dirname(os.path.abspath(__file__))
            absolute_template_path = os.path.join(script_dir, template_path_relative)
            if os.path.exists(absolute_template_path): prs = Presentation(absolute_template_path); print(f"Loaded template: {absolute_template_path}")
            else: print(f"Warning: Template not found at {absolute_template_path}")
        except NameError: # __file__ not defined
             if os.path.exists(template_path_relative):
                 try: prs = Presentation(template_path_relative); print(f"Loaded template from CWD: {template_path_relative}")
                 except Exception as e_cwd: print(f"Error loading template from CWD: {e_cwd}")
             else: print(f"Warning: Template not found in CWD.")
        except Exception as e: print(f"Error loading template '{template_path_relative}': {e}")
    if prs is None: print("Using default pptx template."); prs = Presentation()

    # --- Slide 1: Title Slide (using layout 1 assuming layout 0 is the one to delete) ---
    # We add this *first*, then delete the original slide 0 later.
    try:
        # Use layout 1 (typically Title Slide layout in default templates)
        # Adjust index if your template's title slide layout is different
        title_slide_layout_index = 1
        title_layout = prs.slide_layouts[title_slide_layout_index]
        slide = prs.slides.add_slide(title_layout)
        title_shape = None; subtitle_shape = None
        try: title_shape = slide.shapes.title
        except AttributeError: 
            try: title_shape = slide.placeholders[0]
            except (KeyError, IndexError): pass
        try: subtitle_shape = slide.placeholders[1] # Often subtitle placeholder index
        except (KeyError, IndexError): pass
        if title_shape: add_text_to_shape(title_shape, "Recommendation Report")
        else: print("Warning: Title slide layout missing title placeholder.")
        if subtitle_shape: add_text_to_shape(subtitle_shape, f"Session ID: {session_id}\nGenerated: {time.strftime('%Y-%m-%d')}")
        else: print("Warning: Title slide layout missing subtitle placeholder.")
    except IndexError:
         print(f"Error: Layout index {title_slide_layout_index} not found for Title Slide. Using layout 0.")
         try: # Fallback to layout 0 if 1 doesn't exist
             title_layout = prs.slide_layouts[0]; slide = prs.slides.add_slide(title_layout)
             # Try adding title/subtitle to layout 0... (same logic as above)
             title_shape = None; subtitle_shape = None
             try: title_shape = slide.shapes.title
             except AttributeError: 
                try: title_shape = slide.placeholders[0]
                except (KeyError, IndexError): pass
             try: subtitle_shape = slide.placeholders[1]
             except (KeyError, IndexError): pass
             if title_shape: add_text_to_shape(title_shape, "Recommendation Report")
             if subtitle_shape: add_text_to_shape(subtitle_shape, f"Session ID: {session_id}\nGenerated: {time.strftime('%Y-%m-%d')}")
         except Exception as e_fallback: print(f"Error creating title slide on fallback layout 0: {e_fallback}")
    except Exception as e: print(f"Error creating title slide: {e}")


    # --- Process fragments ---
    sorted_fragment_ids = sorted(report_data.keys())
    for fragment_id in sorted_fragment_ids:
        data = report_data[fragment_id]; frag_info = data.get("fragment_info", {})
        # *** Use frag_names from DB ***
        frag_display_name = frag_info.get('frag_names') or f"Fragment {fragment_id}" # Fallback if name missing
        frag_display_name = str(frag_display_name).strip()
        print(f"Generating slides for Fragment: {fragment_id} ({frag_display_name})...")

        # --- REMOVED: Fragment Title Slide ---
        # try:
        #     header_layout = prs.slide_layouts[DEFAULT_LAYOUT_SECTION_HEADER]; slide = prs.slides.add_slide(header_layout); # ...
        # except ...

        # Slide: Insights (use display name)
        insights = data.get("insights", "Insights could not be generated.");
        insights_title = f"{frag_display_name}: Top 3 Insights"; # Use Display Name
        insights_cleaned = insights
        if isinstance(insights, str): insights_cleaned = re.sub(r'^\s*(\d+\.|-|\*)\s*', '', insights, flags=re.MULTILINE).strip().replace('***', '').replace('---', '')
        add_text_slide(prs, insights_title, insights_cleaned if insights_cleaned else "Insights empty/failed.")

        # Slides: Roadmap Overview & CONSOLIDATED Details
        parsed_roadmap_data = data.get("roadmap")
        # Use Display Name for titles
        roadmap_overview_title = f"{frag_display_name}: Roadmap Overview"
        roadmap_details_title = f"{frag_display_name}: Implementation Plan Details"

        is_valid_parsed_roadmap = isinstance(parsed_roadmap_data, dict) and "error" not in parsed_roadmap_data.get("details", {})
        has_graph_data = isinstance(parsed_roadmap_data, dict) and parsed_roadmap_data.get("phases_for_graph")

        if is_valid_parsed_roadmap or has_graph_data:
            # Roadmap Overview slide
            try:
                total_hrs = parsed_roadmap_data.get('total_estimated_hours', 0)
                # *** Use consolidated description instead of summary_context ***
                overview_content = parsed_roadmap_data.get('details', {}).get('description', 'N/A')
                # Append total hours and roles
                overview_text = f"Overview:\n{overview_content}\n\nTotal Est: {total_hrs} hrs (~{total_hrs/HOURS_PER_DAY:.1f} days)\n\nKey Roles: {parsed_roadmap_data.get('key_roles_involved', 'N/A')}"
                add_text_slide(prs, roadmap_overview_title, overview_text, DEFAULT_LAYOUT_TITLE_CONTENT)
            except Exception as e: print(f"Error creating roadmap overview slide: {e}")

            # Single Consolidated Roadmap Details Slide
            consolidated_details = parsed_roadmap_data.get('details', {})
            if "error" not in consolidated_details:
                 # *** Format content string WITHOUT who_roles_involved ***
                 consolidated_content_text = (
                     f"Overall Goal: {consolidated_details.get('title', 'N/A')}\n\n"
                     # Description is now shown on overview slide, optionally repeat or omit here
                     # f"Description:\n{consolidated_details.get('description', 'N/A')}\n\n"
                     f"Methodology / Steps:\n{consolidated_details.get('how_methodology', 'N/A')}\n\n"
                     # f"Roles Involved (Overall): {consolidated_details.get('who_roles_involved', 'N/A')}\n" # REMOVED
                     f"Tools/Platforms (Overall): {consolidated_details.get('tools_platforms_used', 'N/A')}\n\n"
                     f"Subtask Estimates (Aggregated):\n{consolidated_details.get('estimate_per_subtask', 'N/A')}"
                 )
                 add_text_slide(prs, roadmap_details_title, consolidated_content_text, DEFAULT_LAYOUT_TITLE_CONTENT)
            else:
                 add_text_slide(prs, roadmap_details_title, f"Roadmap details could not be processed.\nError: {consolidated_details.get('error')}", DEFAULT_LAYOUT_TITLE_CONTENT)

        else: # Handle case where synthesis/parsing failed entirely
            roadmap_error_msg = "Roadmap could not be generated or parsed."
            if isinstance(parsed_roadmap_data, dict): roadmap_error_msg += f"\n\nDetails: {parsed_roadmap_data.get('error', parsed_roadmap_data.get('details', {}).get('error', 'Unknown parsing error'))}"
            add_text_slide(prs, roadmap_overview_title, roadmap_error_msg)
            add_text_slide(prs, roadmap_details_title, "Roadmap details unavailable.")

        # Slide: Roadmap Graph (use display name)
        graph_buffer = data.get("graph_buffer")
        graph_title = f"{frag_display_name}: Roadmap Phase Durations" # Use Display Name
        if graph_buffer:
            try:
                graph_layout = None
                try: graph_layout = prs.slide_layouts[DEFAULT_LAYOUT_TITLE_ONLY]
                except IndexError: 
                    try: graph_layout = prs.slide_layouts[DEFAULT_LAYOUT_TITLE_CONTENT]
                    except IndexError: print(f"Error: Layouts {DEFAULT_LAYOUT_TITLE_ONLY}/{DEFAULT_LAYOUT_TITLE_CONTENT} missing."); graph_layout = None
                if graph_layout:
                    slide = prs.slides.add_slide(graph_layout); title_shape_g = None
                    try: title_shape_g = slide.shapes.title
                    except AttributeError: 
                        try: title_shape_g = slide.placeholders[0]
                        except (KeyError, IndexError): pass
                    if title_shape_g: add_text_to_shape(title_shape_g, graph_title)
                    else: print(f"Warning: Graph slide missing title placeholder.")
                    left, top, height = Inches(0.5), Inches(1.5), Inches(5.5)
                    try: pic = slide.shapes.add_picture(graph_buffer, left, top, height=height)
                    except Exception as pic_err: print(f"Error adding picture to graph slide: {pic_err}")
                else: print("Skipping graph slide - suitable layout not found.")
            except Exception as e: print(f"Error adding graph slide: {e}"); traceback.print_exc()
            finally:
                 if graph_buffer and not graph_buffer.closed: graph_buffer.close()
        elif has_graph_data:
             add_text_slide(prs, graph_title, "Graph could not be generated (e.g., zero durations).")
        # --- End Fragment Processing ---

    # *** NEW: Delete the first slide (template reference slide) ***
    try:
        if len(prs.slides._sldIdLst) > 0:
             # Access internal list to delete slide by index
             xml_slides = prs.slides._sldIdLst
             slides = list(xml_slides)
             prs.part.drop_rel(slides[0].rId) # Drop relationship
             xml_slides.remove(slides[0]) # Remove from list
             print("Removed the first slide from the presentation.")
        else:
             print("Warning: No slides found in presentation to remove the first one.")
    except Exception as e_del:
        print(f"Error removing first slide: {e_del}")
        traceback.print_exc()


    # --- Save final presentation to buffer ---
    try:
        ppt_buffer = io.BytesIO(); prs.save(ppt_buffer); ppt_buffer.seek(0); print("\nPPTX saved to memory buffer."); return ppt_buffer
    except Exception as e: print(f"\nError saving presentation to buffer: {e}"); traceback.print_exc(); return None

# --- Upload Function (Unchanged) ---
def upload_report_to_supabase(session_id: str, ppt_buffer: io.BytesIO, bucket_name: str = "recommendations") -> bool:
    """Uploads the generated PPTX buffer's content (bytes) to Supabase Storage."""
    # ... (function body same as before) ...
    if not ppt_buffer: print("Upload failed: No PPT buffer provided."); return False
    remote_path = f"{session_id}/recommendation_report.pptx"
    try:
        print(f"Attempting upload: bucket='{bucket_name}', path='{remote_path}'")
        supabase_client = get_supabase_client(); ppt_buffer.seek(0); ppt_bytes = ppt_buffer.read()
        supabase_client.storage.from_(bucket_name).upload(
            path=remote_path, file=ppt_bytes,
            file_options={"content-type": "application/vnd.openxmlformats-officedocument.presentationml.presentation", "upsert": "true"}
        )
        print(f"Supabase upload finished for {remote_path}. Check dashboard.")
        return True
    except Exception as e: print(f"Error uploading to Supabase: {e}"); traceback.print_exc(); return False
    finally:
        if ppt_buffer and not ppt_buffer.closed: ppt_buffer.close()

# --- Main Generation Function (Chunking/Synthesis - Needs updated parser call) ---
def generate_recommendations(session_id: str) -> Dict[str, Dict[str, Any]]:
    """Generates recommendations using context-aware chunking and consolidated synthesis."""
    print(f"Starting recommendation generation for session_id: {session_id}")
    supabase_client = get_supabase_client(); print("Fetching data...")
    followup_resp = fetch_followup_data(supabase_client, session_id)
    if not followup_resp: print("No followup responses found."); return {}
    q_ids = list(set(r['question_id'] for r in followup_resp if r.get('question_id')))
    q_details = fetch_question_details(supabase_client, q_ids)
    if not q_details: print("No question details found."); return {}
    print("Processing and grouping data...")
    grouped_data = group_data_by_fragment(followup_resp, q_details) # Now gets frag_names
    if not grouped_data: print("No data grouped by fragment."); return {}

    final_report_data = {}; sorted_fragment_ids = sorted(grouped_data.keys()); token_counter_model = None
    try: token_counter_model = genai.GenerativeModel(LLM_MODEL_NAME); print("Token counter model initialized.")
    except Exception as e: print(f"Warning: Could not initialize token counter model: {e}.")

    for fragment_id in sorted_fragment_ids:
        fragment_data = grouped_data[fragment_id]; frag_info = fragment_data.get("fragment_info", {})
        # *** Use frag_names from DB ***
        frag_display_name = frag_info.get('frag_names') or f"Fragment {fragment_id}"
        frag_display_name = str(frag_display_name).strip()
        print(f"\n--- Processing Fragment: {fragment_id} ({frag_display_name}) ---")
        all_questions = fragment_data.get("questions_data", [])
        if not all_questions: print("No questions found. Skipping fragment."); final_report_data[fragment_id] = {"fragment_info": frag_info, "insights": "No questions.", "roadmap": None, "graph_buffer": None}; continue

        final_report_data[fragment_id] = {"fragment_info": frag_info}; partial_insights_list = []; partial_roadmaps_list = []; current_chunk_questions = []; current_chunk_tokens = 0; chunk_num = 1

        # Chunking Loop (Generates partial 3-phase roadmaps)
        for i, question_data in enumerate(all_questions):
            temp_chunk_for_estimation = current_chunk_questions + [question_data]; potential_chunk_str = format_chunk_for_prompt(frag_info, temp_chunk_for_estimation); estimated_next_chunk_tokens = 0
            try:
                if token_counter_model: estimated_next_chunk_tokens = token_counter_model.count_tokens(potential_chunk_str).total_tokens
                else: estimated_next_chunk_tokens = len(potential_chunk_str) // 4
            except Exception as e: print(f"Warning: Token counting failed during chunking ({e}). Using fallback."); estimated_next_chunk_tokens = len(potential_chunk_str) // 4
            is_last_question = (i == len(all_questions) - 1)
            # --- process_chunk helper function (same as before) ---
            def process_chunk(questions_in_chunk, current_chunk_num):
                nonlocal partial_insights_list, partial_roadmaps_list
                print(f"\nChunk {current_chunk_num}: Processing {len(questions_in_chunk)} questions.")
                chunk_context_str = format_chunk_for_prompt(frag_info, questions_in_chunk)
                insights_prompt = generate_insights_prompt(chunk_context_str); insights_result = call_llm(insights_prompt, max_tokens=800, temperature=0.5)
                if insights_result: partial_insights_list.append(insights_result); print(f"  Chunk {current_chunk_num} Insights: Generated.")
                else: print(f"  Chunk {current_chunk_num} Insights: FAILED.")
                roadmap_prompt = generate_roadmap_prompt(chunk_context_str); roadmap_txt = call_llm(roadmap_prompt, max_tokens=3500, temperature=0.3)
                if roadmap_txt:
                    try: # Store raw JSON
                        cleaned_roadmap_txt = roadmap_txt.strip()
                        cleaned_roadmap_txt = re.sub(r'^```(?:json)?\s*', '', cleaned_roadmap_txt, flags=re.IGNORECASE)
                        cleaned_roadmap_txt = re.sub(r'\s*```$', '', cleaned_roadmap_txt).strip()
                        partial_roadmap_json = json.loads(cleaned_roadmap_txt)
                        partial_roadmaps_list.append(partial_roadmap_json)
                        print(f"  Chunk {current_chunk_num} Roadmap: Parsed (raw).")
                    except json.JSONDecodeError as json_err: print(f"  Chunk {current_chunk_num} Roadmap: Raw JSON Parse FAILED. Error: {json_err}"); print(f"RAW TEXT:\n{roadmap_txt}\n")
                else: print(f"  Chunk {current_chunk_num} Roadmap: Generation FAILED.")
            # --- End process_chunk ---
            if current_chunk_questions and estimated_next_chunk_tokens > MAX_CHUNK_TOKENS and not is_last_question:
                process_chunk(current_chunk_questions, chunk_num)
                current_chunk_questions = [question_data]; chunk_num += 1; new_chunk_str = format_chunk_for_prompt(frag_info, current_chunk_questions)
                try:
                     if token_counter_model: current_chunk_tokens = token_counter_model.count_tokens(new_chunk_str).total_tokens
                     else: current_chunk_tokens = len(new_chunk_str) // 4
                except Exception as e: print(f"Warning: Token counting failed starting chunk {chunk_num} ({e})."); current_chunk_tokens = len(new_chunk_str) // 4
            else:
                current_chunk_questions.append(question_data); current_chunk_tokens = estimated_next_chunk_tokens
                if is_last_question: process_chunk(current_chunk_questions, chunk_num)

        # --- Synthesis Step ---
        print("\n--- Synthesizing Results for Fragment ---")
        # Synthesize Insights (same as before)
        final_insights = "Insights synthesis skipped (no partial insights generated)."
        if partial_insights_list:
            print("Synthesizing insights..."); synthesis_insights_prompt = generate_synthesize_insights_prompt(partial_insights_list)
            if synthesis_insights_prompt:
                 final_insights_result = call_llm(synthesis_insights_prompt, max_tokens=1000, temperature=0.6)
                 if final_insights_result: final_insights = final_insights_result; print("Insights synthesized successfully.")
                 else: final_insights = "Insights synthesis failed."; print("Insights synthesis failed.")
            else: final_insights = "Insights synthesis skipped."; print("Insights synthesis skipped.")
        final_report_data[fragment_id]["insights"] = final_insights

        # Synthesize Roadmaps (using the NEW synthesis prompt and UPDATED parser)
        parsed_final_roadmap = None
        if partial_roadmaps_list:
            print("Synthesizing roadmap into consolidated structure...");
            # Use the REWRITTEN synthesis prompt (V4)
            synthesis_roadmaps_prompt = generate_synthesize_roadmaps_prompt(partial_roadmaps_list, fragment_id, frag_display_name) # Pass display name
            if synthesis_roadmaps_prompt:
                 final_roadmap_txt = call_llm(synthesis_roadmaps_prompt, max_tokens=4096, temperature=0.2)
                 if final_roadmap_txt:
                     print("Parsing synthesized CONSOLIDATED roadmap JSON...")
                     # Use the UPDATED parser that expects the new structure
                     parsed_final_roadmap = parse_roadmap_output_json(final_roadmap_txt)
                     if parsed_final_roadmap:
                         # Ensure fragment name is correct in parsed data
                         if parsed_final_roadmap.get('fragment_name') != frag_display_name:
                             print(f"Updating fragment name in parsed data to '{frag_display_name}'")
                             parsed_final_roadmap['fragment_name'] = frag_display_name
                         print("Final consolidated roadmap parsed successfully.")
                     else: print("Failed to parse synthesized consolidated roadmap JSON."); parsed_final_roadmap = {"error": "Failed to parse synthesized roadmap JSON", "raw_output": final_roadmap_txt[:1000]+"..."}
                 else: print("Failed to generate synthesized roadmap JSON text."); parsed_final_roadmap = {"error": "Failed to generate synthesized roadmap JSON text"}
            else: print("Roadmap synthesis skipped (prompt generation failed)."); parsed_final_roadmap = {"error": "Roadmap synthesis prompt generation failed"}
        else: print("Roadmap synthesis skipped (no valid partial roadmaps were generated/parsed)."); parsed_final_roadmap = {"error": "No valid partial roadmaps generated/parsed"}
        final_report_data[fragment_id]["roadmap"] = parsed_final_roadmap # Store result

        # Generate Graph (uses the 'phases_for_graph' extracted by the updated parser)
        print("Generating graph..."); graph_buffer = None
        if isinstance(parsed_final_roadmap, dict) and parsed_final_roadmap.get('phases_for_graph'):
            graph_buffer = create_roadmap_graph(parsed_final_roadmap['phases_for_graph'])
        final_report_data[fragment_id]["graph_buffer"] = graph_buffer
        # --- End Fragment Processing ---

    print("\nRecommendation generation process completed.")
    return final_report_data


# --- Main Orchestration Function (Unchanged) ---
def generate_and_upload_recommendations(session_id: str, template_path_relative: Optional[str] = "template.pptx") -> bool:
     """Generates recommendations, builds PPT, and uploads."""
     # ... (function body same as before) ...
     print(f"Starting generation and upload for session {session_id}...")
     try:
         report_content = generate_recommendations(session_id)
         if not report_content: print(f"Failed to generate content for session {session_id}."); return False
         print(f"Building PPT report for session {session_id}...")
         ppt_buffer = build_ppt_report(session_id, report_content, template_path_relative)
         if not ppt_buffer:
             print(f"Failed to build PPT buffer for session {session_id}.")
             for frag_id in report_content: # Cleanup graph buffers
                 if report_content[frag_id].get('graph_buffer'):
                     try: report_content[frag_id]['graph_buffer'].close()
                     except Exception: pass
             return False
         print(f"Uploading report for session {session_id}...")
         upload_success = upload_report_to_supabase(session_id, ppt_buffer) # Buffer closed inside
         print(f"Successfully generated and uploaded report for session {session_id}." if upload_success else f"Upload failed for session {session_id}.")
         return upload_success
     except Exception as e: print(f"Error during generate_and_upload for session {session_id}: {e}"); traceback.print_exc(); return False


# --- Main Execution Block (Hardcoded Local Save) ---
if __name__ == "__main__":
    # --- Hardcoded Settings ---
    session_id_to_process = "df033ae9-84d4-48ba-a577-374955cbe690" # <-- EDIT THIS LINE
    template_file_name = "template.pptx"
    hardcoded_output_filename = "Generated_Recommendation_Report_Streamlined.pptx" # New name
    should_upload = False # Set True to also upload
    # ---

    parser = argparse.ArgumentParser(description="Generate recommendations PPT from session_id.")
    parser.add_argument("-s", "--session", type=str, default=session_id_to_process, help="Session ID to process.")
    parser.add_argument("-t", "--template", type=str, default=template_file_name, help="Relative path to PPTX template.")
    parser.add_argument("--upload", action='store_true', default=should_upload, help="Upload report to Supabase.")
    args = parser.parse_args(); session_id_to_process = args.session; template_file_name = args.template; should_upload = args.upload
    if not session_id_to_process or session_id_to_process == "YOUR_SESSION_ID_HERE": print("Error: Provide valid session ID via --session or edit script."); sys.exit(1)
    print(f"Processing session_id: {session_id_to_process}"); print(f"Using template: {template_file_name}"); print(f"Hardcoded local output: {hardcoded_output_filename}"); print(f"Upload enabled: {should_upload}")

    report_content = generate_recommendations(session_id_to_process)
    if report_content:
        print("\nBuilding PowerPoint presentation...")
        ppt_buffer = build_ppt_report(session_id_to_process, report_content, template_path_relative=template_file_name)
        if ppt_buffer:
            print("PPT Buffer created.")
            try: # Save locally
                with open(hardcoded_output_filename, 'wb') as f: f.write(ppt_buffer.getvalue())
                print(f"Presentation saved locally: {hardcoded_output_filename}")
            except Exception as e: print(f"\nError saving locally to '{hardcoded_output_filename}': {e}")
            if should_upload: # Upload if requested
                 print("\nUploading report to Supabase..."); ppt_buffer.seek(0); upload_success = upload_report_to_supabase(session_id_to_process, ppt_buffer)
                 print("Upload successful." if upload_success else "Upload failed.")
            elif not ppt_buffer.closed: ppt_buffer.close() # Close buffer if not uploaded
        else:
            print("\nFailed to build PPT buffer.")
            for frag_id in report_content: # Cleanup graph buffers
                 if report_content[frag_id].get('graph_buffer'):
                     try: report_content[frag_id]['graph_buffer'].close()
                     except Exception: pass
    else: print("\nNo report data generated. PPT creation skipped.")
    print("\nScript finished.")
